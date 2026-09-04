from io import BytesIO

from fastapi import APIRouter, UploadFile, File, HTTPException
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from app.engines.material_engine import material_engine


router = APIRouter(
    prefix="/api/material",
    tags=["Material Intelligence"]
)


def extract_pdf(data: bytes) -> str:
    reader = PdfReader(BytesIO(data))

    pages = []

    for page in reader.pages:
        text = page.extract_text() or ""
        pages.append(text)

    return "\n".join(pages)


def extract_docx(data: bytes) -> str:
    document = Document(BytesIO(data))

    return "\n".join(
        paragraph.text
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    )


def extract_pptx(data: bytes) -> str:
    presentation = Presentation(BytesIO(data))

    slides = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slides.append(text)

    return "\n".join(slides)


def extract_text(filename: str, data: bytes) -> str:

    extension = filename.lower().split(".")[-1]

    if extension == "txt":
        return data.decode("utf-8", errors="ignore")

    if extension == "pdf":
        return extract_pdf(data)

    if extension == "docx":
        return extract_docx(data)

    if extension == "pptx":
        return extract_pptx(data)

    if extension == "doc":
        raise HTTPException(
            status_code=400,
            detail="Legacy .doc files are not supported yet. Please use .docx."
        )

    if extension == "ppt":
        raise HTTPException(
            status_code=400,
            detail="Legacy .ppt files are not supported yet. Please use .pptx."
        )

    raise HTTPException(
        status_code=400,
        detail=f"Unsupported file format: .{extension}"
    )


@router.post("/analyze")
async def analyze_material(file: UploadFile = File(...)):

    data = await file.read()

    if not data:
        raise HTTPException(
            status_code=400,
            detail="Uploaded file is empty."
        )

    text = extract_text(
        file.filename or "material.txt",
        data
    )

    analysis = material_engine.analyze(text)

    analysis["file"] = {
        "filename": file.filename,
        "content_type": file.content_type
    }

    # Preserve the extracted learning-material text so the
    # competency pipeline can continue processing the same upload.
    analysis["text"] = text

    return analysis


@router.post("/analyze-text")
async def analyze_text(payload: dict):

    text = payload.get("text", "")

    if not text:
        raise HTTPException(
            status_code=400,
            detail="Text is required."
        )

    return material_engine.analyze(text)
